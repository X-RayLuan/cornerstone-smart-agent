from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


BASE = Path("/Users/leiwang")
DESKTOP = BASE / "Desktop/Cornerstone Wiki"
WORKSPACE = BASE / ".openclaw/workspace-smart"
TMP = WORKSPACE / ".tmp/template_media"
TEMPLATE = DESKTOP / "柯纳仕咨询与培训 Intro_CN 咨询版 V1.2.pptx"
OUTPUT = DESKTOP / f"百威商学院_管理与领导力发展课程方案_柯纳仕_{datetime.now():%Y%m%d}.pptx"

GOLD = RGBColor(0x9B, 0x7D, 0x5F)
DARK = RGBColor(0x35, 0x2E, 0x2A)
TAUPE = RGBColor(0xB2, 0x9A, 0x82)
LIGHT = RGBColor(0xF5, 0xF1, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x62, 0x5B)

LOGO_DARK = TMP / "image14.png"
LOGO_LIGHT = TMP / "image5.png"
PHOTO_BRIDGE = TMP / "image2.jpg"
PHOTO_OFFICE = TMP / "image3.jpeg"
PHOTO_WORKSHOP = TMP / "image12.jpeg"
PHOTO_TEAM = TMP / "image11.jpeg"


def add_full_bg(slide, image_path: Path, prs: Presentation) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def add_logo(slide, image_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(image_path), Cm(left), Cm(top), width=Cm(width))


def add_rect(slide, left, top, width, height, fill, line=None, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line or fill
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="PingFang SC",
    margin=0.18,
):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(margin)
    tf.margin_right = Cm(margin)
    tf.margin_top = Cm(margin)
    tf.margin_bottom = Cm(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for raw_line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = raw_line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return box


def add_title(slide, title: str, subtitle: Optional[str] = None) -> None:
    add_textbox(slide, 1.3, 0.9, 18.8, 1.1, title, font_size=24, color=GOLD, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.35), Cm(2.0), Cm(3.4), Cm(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 1.35, 2.18, 24.5, 0.8, subtitle, font_size=10.5, color=MUTED)


def add_bullet_block(slide, left, top, width, height, heading, bullets, fill=WHITE):
    add_rect(slide, left, top, width, height, fill, line=LIGHT)
    add_textbox(slide, left + 0.18, top + 0.12, width - 0.36, 0.7, heading, font_size=14, color=GOLD, bold=True)
    bullet_lines = "\n".join([f"• {item}" for item in bullets])
    add_textbox(slide, left + 0.18, top + 0.88, width - 0.36, height - 1.0, bullet_lines, font_size=10.5)


def add_course_slide(slide, title, audience, problems, modules, scenes, extra) -> None:
    add_title(slide, title, "课程信息保留知识库原有细节，并结合百威商学院应用场景进行排布。")
    add_rect(slide, 1.25, 2.85, 24.0, 9.8, LIGHT, line=LIGHT)

    add_bullet_block(slide, 1.6, 3.2, 7.4, 3.0, "适用对象", audience)
    add_bullet_block(slide, 9.25, 3.2, 7.7, 3.0, "主要解决的问题", problems)
    add_bullet_block(slide, 17.25, 3.2, 7.6, 3.0, "典型应用场景", scenes)
    add_bullet_block(slide, 1.6, 6.55, 14.0, 4.7, "典型内容模块", modules)
    add_bullet_block(slide, 15.95, 6.55, 8.9, 4.7, "在百威商学院中的建议用法", extra)
    add_logo(slide, LOGO_DARK, 21.5, 0.92, 3.0)


def build() -> Path:
    template_prs = Presentation(str(TEMPLATE))
    prs = Presentation()
    prs.slide_width = template_prs.slide_width
    prs.slide_height = template_prs.slide_height
    layout = prs.slide_layouts[6]

    # Cover
    slide = prs.slides.add_slide(layout)
    add_full_bg(slide, PHOTO_BRIDGE, prs)
    panel = add_rect(slide, 1.1, 1.0, 12.0, 10.0, WHITE, line=WHITE, transparency=0.10)
    panel.fill.transparency = 0.12
    add_logo(slide, LOGO_LIGHT, 1.7, 1.45, 3.6)
    add_textbox(slide, 1.7, 3.1, 10.6, 2.3, "百威商学院\n管理与领导力发展课程方案", font_size=24, color=DARK, bold=True)
    add_textbox(
        slide,
        1.8,
        5.65,
        9.8,
        2.6,
        "基于《柯纳仕咨询与培训 Intro_CN 咨询版 V1.2》视觉模板生成\n聚焦中国区培训部门的管理者与关键人才发展需求",
        font_size=12,
        color=MUTED,
    )
    add_textbox(slide, 1.8, 9.8, 6.0, 0.8, f"柯纳仕咨询 | {datetime.now():%Y.%m.%d}", font_size=10.5, color=TAUPE)

    # Client understanding
    slide = prs.slides.add_slide(layout)
    add_title(slide, "项目理解与定制逻辑", "围绕百威商学院作为中国区培训部门的角色，优先匹配最常见的人才发展与管理协同场景。")
    add_bullet_block(
        slide,
        1.4,
        3.0,
        7.5,
        4.0,
        "本方案优先关注的人群",
        [
            "个人贡献者与储备管理者",
            "新任经理与一线带队主管",
            "中高层业务/职能管理者",
            "需要推动跨部门协同的关键岗位人才",
        ],
    )
    add_bullet_block(
        slide,
        9.2,
        3.0,
        7.8,
        4.0,
        "优先对应的典型业务场景",
        [
            "销售、市场、供应链、工厂、总部职能之间的协同",
            "新经理带队与班组/团队管理",
            "跨区域或跨部门项目推进",
            "关键人才梯队建设与领导力升级",
        ],
    )
    add_bullet_block(
        slide,
        17.35,
        3.0,
        7.6,
        4.0,
        "课程推荐原则",
        [
            "优先使用柯纳仕现有知识库中的成熟课程",
            "不删减课程页中的关键细节",
            "课程与方法论组合，兼顾训练与转化",
            "既适配通用管理问题，也适配业务落地场景",
        ],
    )
    add_rect(slide, 1.4, 7.55, 23.6, 3.9, WHITE, line=LIGHT)
    add_textbox(
        slide,
        1.75,
        7.9,
        22.8,
        2.9,
        "定制建议：百威商学院可将本方案作为“基础管理能力池 + 跨部门协同能力池 + 领导力进阶能力池”的课程底盘，按人群分层配置，\n并在重点项目中叠加行动学习、复盘与教练辅导，增强从课堂到业务场景的转化。",
        font_size=14,
    )
    add_logo(slide, LOGO_DARK, 21.5, 0.92, 3.0)

    # Company and differentiation
    slide = prs.slides.add_slide(layout)
    add_title(slide, "为什么是柯纳仕", "保持现有课程细节的同时，用项目化与方法论把学习真正连到业务场景。")
    add_bullet_block(
        slide,
        1.4,
        3.0,
        7.2,
        7.8,
        "柯纳仕的定位",
        [
            "聚焦领导力发展、组织发展与管理能力提升",
            "不仅提供标准课程，也能围绕真实业务场景设计项目化解决方案",
            "覆盖个人贡献者、新经理、中高层、C-level 的完整发展路径",
            "产品形态包括标准课程、训练营/发展项目、行动学习项目、咨询式解决方案",
        ],
    )
    add_bullet_block(
        slide,
        8.95,
        3.0,
        7.6,
        7.8,
        "与单点培训的差异",
        [
            "不是只卖标准课件，而是按层级与场景设计学习路径",
            "不是培训结束即结束，而是可叠加复盘、教练、行动学习形成成长闭环",
            "不是面向所有人同一套内容，而是按角色与组织挑战匹配主题",
            "不是只看课堂体验，而是关注行为改变与项目转化",
        ],
    )
    add_bullet_block(
        slide,
        16.9,
        3.0,
        8.1,
        7.8,
        "对百威商学院的价值",
        [
            "能快速形成分层课程地图，支撑商学院年度课程池建设",
            "可把关键课程嵌入新经理、高潜、跨部门项目等重点人群项目",
            "可将课程与真实业务课题结合，减少“学完回不去”的问题",
            "后续可持续扩展到青年人才、行动学习、专项咨询项目",
        ],
    )
    add_logo(slide, LOGO_DARK, 21.5, 0.92, 3.0)

    # Capability map
    slide = prs.slides.add_slide(layout)
    add_title(slide, "能力地图与推荐结构", "采用“层级 × 能力主题”的双轴结构，将课程推荐和人才发展路径对齐。")
    levels = [
        ("个人贡献者", ["问题分析与解决", "高效沟通", "非职权影响力", "时间管理/项目管理"]),
        ("新任经理/中基层", ["建设高效团队", "绩效管理", "授权激励", "教练式管理与辅导"]),
        ("中高层管理者", ["教练式领导力", "跨部门协同", "利益相关者管理", "变革与创新"]),
        ("高层管理者", ["战略解码", "内在领导力", "人际领导力", "经营/人才沙盘"]),
    ]
    left = 1.5
    for title, items in levels:
        add_rect(slide, left, 3.2, 5.75, 6.2, WHITE, line=LIGHT)
        add_textbox(slide, left + 0.2, 3.45, 5.3, 0.8, title, font_size=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + 0.25, 4.35, 5.2, 4.8, "\n".join([f"• {x}" for x in items]), font_size=11)
        left += 6.0
    add_rect(slide, 1.5, 10.0, 23.5, 1.6, LIGHT, line=LIGHT)
    add_textbox(
        slide,
        1.8,
        10.25,
        22.8,
        1.0,
        "本次优先保留并展开的 5 门课程：高效沟通、非职权影响力、问题分析与决策、建设高效团队、教练式领导力。",
        font_size=12.5,
        color=DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_logo(slide, LOGO_DARK, 21.5, 0.92, 3.0)

    # Course slides
    for spec in [
        (
            "推荐课程 01｜高效沟通",
            ["需要提升跨部门沟通效率的员工", "新任经理或储备管理者", "需要提升表达清晰度与协作质量的团队成员"],
            ["表达不清，导致协作成本高", "沟通中容易产生误解或冲突", "只会“说”，不会根据对象调整表达方式", "会议、汇报、协作中的沟通结果不稳定"],
            ["沟通基本模型与常见误区", "倾听与澄清", "面向不同对象的表达调整", "工作场景中的反馈与对话", "冲突或分歧下的沟通处理"],
            ["团队协作沟通", "跨部门项目配合", "向上汇报与横向沟通", "反馈与绩效对话前的基础训练"],
            ["适合作为基层管理与储备干部课程池的基础模块", "可前置于非职权影响力、工作汇报、冲突管理等主题", "适合销售/供应链/职能协同场景中的共同语言建设"],
        ),
        (
            "推荐课程 02｜非职权影响力",
            ["个人贡献者", "项目经理", "跨部门协作频繁的员工", "新经理与中基层管理者"],
            ["有专业意见，但推动不动别人", "需要跨部门协作，却缺少资源与支持", "只会讲自己的逻辑，忽略对方立场", "沟通中容易陷入硬碰硬或无效妥协"],
            ["影响力的基本机制", "识别利益相关者诉求", "建立信任与信用", "说服逻辑与论证表达", "在分歧与阻力中推动共识"],
            ["项目推进与协调", "向上争取支持", "横向协作与资源整合", "在无直接汇报关系下推动行动"],
            ["适合跨部门项目、区域协同、资源协调类岗位", "建议与高效沟通、问题分析与决策串联成“协同推进”学习路径", "适合解决“没有权力但要拿结果”的典型业务场景"],
        ),
        (
            "推荐课程 03｜问题分析与决策",
            ["个人贡献者", "项目负责人", "新任经理", "需要处理复杂问题的中基层管理者"],
            ["问题界定不清，容易头痛医头", "决策依据不足，只看表面现象", "分析过程混乱，团队难以达成一致", "解决方案缺乏落地性"],
            ["问题定义与拆解", "根因分析", "结构化思考", "方案生成与比较", "决策原则与风险判断", "行动计划与复盘"],
            ["日常业务问题分析", "项目卡点排查", "跨团队协作中的复杂问题", "行动学习项目前的能力准备"],
            ["适合作为业务骨干与新经理的通用底层能力课程", "可作为行动学习、专项改善项目的前置训练", "适合供应链、工厂、销售运营等需要结构化解决问题的场景"],
        ),
        (
            "推荐课程 04｜建设高效团队",
            ["新经理", "中基层管理者", "团队负责人", "需要提升团队协作质量的业务主管"],
            ["团队目标不清或对齐不足", "成员各自为战，协作成本高", "团队缺少复盘与改进机制", "管理者只盯任务，忽略团队系统建设"],
            ["高效团队的关键特征", "目标共识与角色分工", "团队协作机制", "沟通、冲突与信任建设", "团队复盘与持续改进"],
            ["新团队组建", "团队协作效率偏低", "管理者带队能力提升", "团队转型或业务变化阶段"],
            ["适合一线带队主管、新经理、核心班组管理者", "建议与绩效管理、教练式领导力、复盘配套使用", "可用于提升团队从“完成任务”到“系统协作”的管理能力"],
        ),
        (
            "推荐课程 05｜教练式领导力",
            ["中基层到中高层管理者", "需要提升带人能力的业务负责人", "正在建设人才梯队的团队主管"],
            ["管理者过度依赖指令式管理", "带团队时“管事多、育人少”", "团队成员自主性不足", "反馈对话效果差，成员成长缓慢"],
            ["教练式领导的核心理念", "从“给答案”转向“促思考”", "有效提问与深度倾听", "发展性反馈", "辅导员工承担责任与形成行动计划"],
            ["一对一辅导", "人才发展对话", "新经理带团队", "团队成员成长遇到瓶颈时的辅导支持"],
            ["适合作为中高层与关键主管的带人能力升级模块", "可与建设高效团队、绩效管理、授权激励构成管理者发展路径", "适合配合商学院人才项目中的导师制与发展对话机制"],
        ),
    ]:
        slide = prs.slides.add_slide(layout)
        add_course_slide(slide, *spec)

    # Methodology
    slide = prs.slides.add_slide(layout)
    add_title(slide, "方法论补充｜行动学习", "在重点项目中，将课程与真实业务课题结合，增强行为改变与业务转化。")
    add_full_bg(slide, PHOTO_TEAM, prs)
    overlay = add_rect(slide, 0.9, 1.0, 24.6, 10.9, WHITE, line=WHITE, transparency=0.15)
    overlay.fill.transparency = 0.16
    add_bullet_block(
        slide,
        1.4,
        3.0,
        7.2,
        7.5,
        "适用情境",
        [
            "问题真实、重要、紧急，且与团队相关",
            "没有现成解决方案，或对当前方案不满意",
            "需要跨角色、跨团队共创解决思路",
            "希望在解决业务问题的同时发展领导力和团队协作能力",
        ],
        fill=WHITE,
    )
    add_bullet_block(
        slide,
        9.0,
        3.0,
        7.4,
        7.5,
        "关键流程",
        [
            "准备阶段：选择真实业务课题、明确目标与授权、组建课题小组、确认 Sponsor",
            "行动创建工作坊：统一课题目标、初步分析问题、制定行动计划",
            "在行动中学习：围绕差距、根因、假设、行动项进行教练跟进",
            "中期/终期复盘：总结个人成长、团队有效性与整体成果",
        ],
        fill=WHITE,
    )
    add_bullet_block(
        slide,
        16.8,
        3.0,
        8.0,
        7.5,
        "对百威商学院的建议",
        [
            "可嵌入跨部门改善项目、关键人才项目、管理者加速器",
            "适合围绕协同、效率、流程优化、团队管理等真实课题展开",
            "Sponsor 需要承担方向引领、资源支持、过程辅导与成果评审",
            "课程训练可作为前置，行动学习负责把能力转成结果",
        ],
        fill=WHITE,
    )
    add_logo(slide, LOGO_DARK, 21.5, 0.92, 3.0)

    # Delivery suggestion
    slide = prs.slides.add_slide(layout)
    add_title(slide, "建议交付方式", "从课程池到项目化落地，建议以“基础模块 + 分层路径 + 转化机制”组合推进。")
    add_rect(slide, 1.3, 3.0, 23.7, 7.9, LIGHT, line=LIGHT)
    add_bullet_block(
        slide,
        1.7,
        3.4,
        7.2,
        6.9,
        "模块 A｜基础协同能力池",
        [
            "高效沟通",
            "非职权影响力",
            "问题分析与决策",
            "适合个人贡献者、储备管理者、跨部门项目成员",
        ],
    )
    add_bullet_block(
        slide,
        9.2,
        3.4,
        7.2,
        6.9,
        "模块 B｜新经理带队能力池",
        [
            "建设高效团队",
            "高效会议引导 / 复盘 / 绩效管理 / 授权激励",
            "适合新经理、班组长、基层带队主管",
            "目标是从“带任务”升级为“带团队”",
        ],
    )
    add_bullet_block(
        slide,
        16.7,
        3.4,
        7.8,
        6.9,
        "模块 C｜中高层领导力进阶",
        [
            "教练式领导力",
            "跨部门协同、利益相关者管理、变革与创新",
            "适合关键管理者、高潜与核心职能负责人",
            "可叠加行动学习、教练辅导与复盘机制",
        ],
    )
    add_logo(slide, LOGO_DARK, 21.5, 0.92, 3.0)

    # Closing
    slide = prs.slides.add_slide(layout)
    add_full_bg(slide, PHOTO_OFFICE, prs)
    panel = add_rect(slide, 13.4, 1.25, 11.0, 9.8, WHITE, line=WHITE, transparency=0.08)
    panel.fill.transparency = 0.08
    add_logo(slide, LOGO_DARK, 15.0, 1.8, 4.0)
    add_textbox(slide, 15.1, 4.1, 8.3, 2.0, "下一步建议", font_size=25, color=GOLD, bold=True)
    add_textbox(
        slide,
        15.2,
        6.0,
        8.2,
        3.6,
        "1. 确认优先服务人群与年度课程池\n2. 选定 2-3 门课程做首批试点\n3. 对重点项目叠加行动学习或复盘机制\n4. 再扩展至青年人才 / 新经理 / 中高层项目化路径",
        font_size=14,
        color=DARK,
    )
    add_textbox(slide, 15.2, 10.1, 7.5, 0.8, "柯纳仕咨询与培训 | The Cornerstone", font_size=11, color=TAUPE)

    if OUTPUT.exists():
        OUTPUT.unlink()
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(out)
