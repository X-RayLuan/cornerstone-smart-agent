#!/usr/bin/env python3
from __future__ import annotations

import argparse
import html
from html.parser import HTMLParser
from pathlib import Path
from typing import List, Optional, Set

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

WORKSPACE = Path("/Users/leiwang/.openclaw/workspace-smart")
MEDIA = WORKSPACE / ".tmp" / "template_media"
LOGO_DARK = MEDIA / "image14.png"
LOGO_LIGHT = MEDIA / "image5.png"
PHOTO_COVER = MEDIA / "image2.jpg"
INTRO_COVER_REF = MEDIA / "intro_cover_ref.png"

GOLD = RGBColor(0x9B, 0x7D, 0x5F)
TAUPE = RGBColor(0xB2, 0x9A, 0x82)
DARK = RGBColor(0x35, 0x2E, 0x2A)
MUTED = RGBColor(0x6B, 0x62, 0x5B)
LIGHT = RGBColor(0xF5, 0xF1, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE3, 0xD9, 0xCE)
TITLE_FONT = "Noto Sans S Chinese Black"
BODY_FONT = "Noto Sans S Chinese Light"
BOLD_FONT = "Noto Sans S Chinese Bold"


class SlideContent:
    def __init__(self) -> None:
        self.title = ""
        self.subtitle = ""
        self.paragraphs: List[str] = []
        self.bullets: List[str] = []
        self.table: List[List[str]] = []
        self.classes: Set[str] = set()


class DeckParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.slides: List[SlideContent] = []
        self.current_slide: Optional[SlideContent] = None
        self.current_tag: Optional[str] = None
        self.current_cell_row: List[str] = []
        self.current_table: List[List[str]] = []
        self.capture: List[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        attrs_dict = dict(attrs)
        classes = set((attrs_dict.get("class") or "").split())
        if tag == "section" and "slide" in classes:
            self.current_slide = SlideContent()
            self.current_slide.classes = classes
            self.slides.append(self.current_slide)
        elif self.current_slide is not None:
            if tag in {"h1", "h2", "p", "li", "td", "th"}:
                self.current_tag = tag
                self.capture = []
            elif tag == "tr":
                self.current_cell_row = []
            elif tag == "table":
                self.current_table = []
            elif tag == "br":
                self.capture.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if self.current_slide is None:
            return
        text = html.unescape("".join(self.capture).strip())
        if tag == "h1" and self.current_tag == "h1":
            self.current_slide.title = text
        elif tag == "h2" and self.current_tag == "h2":
            if not self.current_slide.subtitle:
                self.current_slide.subtitle = text
            else:
                self.current_slide.paragraphs.append(text)
        elif tag == "p" and self.current_tag == "p" and text:
            self.current_slide.paragraphs.append(text)
        elif tag == "li" and self.current_tag == "li" and text:
            self.current_slide.bullets.append(text)
        elif tag in {"td", "th"} and self.current_tag in {"td", "th"}:
            self.current_cell_row.append(text)
        elif tag == "tr" and self.current_cell_row:
            self.current_table.append(self.current_cell_row)
            self.current_cell_row = []
        elif tag == "table" and self.current_table:
            self.current_slide.table = self.current_table
            self.current_table = []
        if tag == self.current_tag:
            self.current_tag = None
            self.capture = []

    def handle_data(self, data: str) -> None:
        if self.current_tag is not None:
            self.capture.append(data)


def add_shape(slide, shape_type, left_cm, top_cm, width_cm, height_cm, fill=WHITE, line=LINE, radius=True):
    shape = slide.shapes.add_shape(
        shape_type, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_rect(slide, left_cm, top_cm, width_cm, height_cm, fill=WHITE, line=LINE):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left_cm, top_cm, width_cm, height_cm, fill, line)


def add_textbox(
    slide,
    left_cm,
    top_cm,
    width_cm,
    height_cm,
    text,
    size_pt,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    font_name=BODY_FONT,
    margin_cm=0.12,
):
    box = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(margin_cm)
    tf.margin_right = Cm(margin_cm)
    tf.margin_top = Cm(margin_cm)
    tf.margin_bottom = Cm(margin_cm)
    tf.clear()
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return box


def add_bullets(slide, left_cm, top_cm, width_cm, height_cm, bullets, size_pt=13, color=DARK):
    box = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(size_pt)
        p.font.color.rgb = color
        p.font.name = BODY_FONT
    return box


def add_table(slide, left_cm, top_cm, width_cm, height_cm, rows: List[List[str]]):
    if not rows:
        return
    cols = max(len(row) for row in rows)
    shape = slide.shapes.add_table(len(rows), cols, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    table = shape.table
    for r, row in enumerate(rows):
        for c in range(cols):
            value = row[c] if c < len(row) else ""
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Cm(0.12)
            cell.margin_right = Cm(0.12)
            cell.margin_top = Cm(0.08)
            cell.margin_bottom = Cm(0.08)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD
            elif c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10.5 if rows and len(rows) > 4 else 12)
                    run.font.name = BOLD_FONT if (r == 0 or c == 0) else BODY_FONT
                    run.font.color.rgb = WHITE if r == 0 else DARK
                    run.font.bold = r == 0 or c == 0


def add_logo(slide, light=False):
    logo = LOGO_LIGHT if light else LOGO_DARK
    if logo.exists():
        slide.shapes.add_picture(str(logo), Cm(24.0), Cm(0.65), width=Cm(3.0))


def add_title(slide, content: SlideContent, idx: int) -> None:
    add_textbox(slide, 1.25, 0.78, 20.5, 0.82, content.title or "Untitled", 22, bold=True, color=GOLD, font_name=TITLE_FONT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.35), Cm(1.82), Cm(3.2), Cm(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    if content.subtitle:
        add_textbox(slide, 1.25, 2.02, 22.8, 0.65, content.subtitle, 10.5, color=MUTED, font_name=BODY_FONT)
    add_logo(slide)
    add_textbox(slide, 24.5, 18.2, 2.4, 0.45, f"{idx:02d}", 8, color=TAUPE, align=PP_ALIGN.RIGHT, font_name=BODY_FONT)


def render_cover(prs: Presentation, content: SlideContent) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1])
    if INTRO_COVER_REF.exists():
        slide.shapes.add_picture(str(INTRO_COVER_REF), 0, 0, width=prs.slide_width, height=prs.slide_height)
    elif PHOTO_COVER.exists():
        slide.shapes.add_picture(str(PHOTO_COVER), 0, 0, width=prs.slide_width, height=prs.slide_height)
    else:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 33.87, 19.05, LIGHT, LIGHT)
    panel = add_rect(slide, 9.2, 5.6, 15.6, 8.2, WHITE, WHITE)
    panel.fill.transparency = 2
    panel.shadow.inherit = False
    if LOGO_LIGHT.exists():
        slide.shapes.add_picture(str(LOGO_LIGHT), Cm(12.0), Cm(7.1), width=Cm(6.2))
    add_textbox(slide, 11.8, 10.8, 10.7, 1.8, content.title, 23, bold=True, color=DARK, align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    if content.subtitle:
        add_textbox(slide, 10.5, 12.7, 13.3, 0.9, content.subtitle, 10.5, color=GOLD, align=PP_ALIGN.CENTER, font_name=BOLD_FONT)
    if content.paragraphs:
        add_textbox(slide, 1.55, 17.2, 10.8, 0.8, content.paragraphs[0], 10.5, color=GOLD, font_name=BOLD_FONT)


def render_cards(slide, content: SlideContent, start_top=3.15):
    cards = content.bullets
    if not cards:
        return
    cols = 3 if len(cards) <= 6 else 4
    w = 7.4 if cols == 3 else 5.7
    gap = 0.45
    h = 2.45 if len(cards) > cols else 3.3
    for i, item in enumerate(cards):
        row = i // cols
        col = i % cols
        left = 1.35 + col * (w + gap)
        top = start_top + row * (h + 0.45)
        add_rect(slide, left, top, w, h, WHITE, LINE)
        parts = item.split("：", 1)
        head = parts[0]
        body = parts[1] if len(parts) > 1 else ""
        add_textbox(slide, left + 0.28, top + 0.18, w - 0.56, 0.55, head, 13, bold=True, color=GOLD, font_name=BOLD_FONT)
        add_textbox(slide, left + 0.28, top + 0.82, w - 0.56, h - 1.0, body, 10.2, color=DARK, font_name=BODY_FONT)


def render_slide(prs: Presentation, content: SlideContent, idx: int) -> None:
    layout_index = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 33.87, 19.05, RGBColor(0xFB, 0xFA, 0xF8), RGBColor(0xFB, 0xFA, 0xF8))
    add_title(slide, content, idx)
    current_top = 3.05
    if content.subtitle:
        current_top = 2.9
    if content.paragraphs:
        add_rect(slide, 1.35, current_top, 24.0, min(2.7, 0.95 + 0.55 * len(content.paragraphs)), WHITE, LINE)
        add_textbox(slide, 1.65, current_top + 0.16, 23.4, min(2.2, 0.7 + 0.5 * len(content.paragraphs)), "\n".join(content.paragraphs), 11.2, color=DARK, font_name=BODY_FONT)
        current_top += min(3.1, 1.25 + 0.58 * len(content.paragraphs))
    if content.table:
        if "fulltable" in content.classes:
            table_top = 2.55
            table_h = 15.2
            add_table(slide, 0.85, table_top, 31.9, table_h, content.table)
        else:
            table_top = max(current_top + 0.1, 3.15)
            table_h = min(13.8, 1.0 + 0.82 * len(content.table))
            add_table(slide, 1.35, table_top, 24.8, table_h, content.table)
    elif content.bullets:
        if "cards" in content.classes:
            render_cards(slide, content, max(current_top + 0.15, 3.15))
        else:
            add_rect(slide, 1.35, current_top + 0.1, 24.0, 10.0, WHITE, LINE)
            add_bullets(slide, 1.75, current_top + 0.45, 23.2, 9.2, content.bullets, size_pt=12)


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert constrained HTML deck to PPTX.")
    parser.add_argument("html_path", type=Path)
    parser.add_argument("output_path", type=Path)
    parser.add_argument("--template", type=Path, default=None)
    args = parser.parse_args()

    parser_obj = DeckParser()
    parser_obj.feed(args.html_path.read_text(encoding="utf-8"))
    if not parser_obj.slides:
        raise SystemExit("No slides found. Expected <section class=\"slide\"> blocks.")

    prs = Presentation(str(args.template)) if args.template else Presentation()
    while len(prs.slides) > 0:
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[0]

    for idx, slide_content in enumerate(parser_obj.slides, start=1):
        if "cover" in slide_content.classes:
            render_cover(prs, slide_content)
        else:
            render_slide(prs, slide_content, idx)

    args.output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output_path))
    print(args.output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
