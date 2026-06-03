"""
百威商学院 × 柯纳仕 领导力发展课程方案
基于 柯纳仕咨询与培训 Intro_CN 咨询版 V1.2.pptx 模板生成
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
import re

src = Path('/Users/leiwang/Desktop/Cornerstone Wiki/柯纳仕咨询与培训 Intro_CN 咨询版 V1.2.pptx')
dst = Path('/Users/leiwang/Desktop/Cornerstone Wiki/百威商学院_柯纳仕领导力发展课程方案_20260526.pptx')

prs = Presentation(str(src))
slides = list(prs.slides)

def clear_and_set_text(shape, new_text, font_size=None, bold=False, color_rgb=None):
    """Clear shape text and set new text."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = new_text
    if font_size:
        run.font.size = Pt(font_size)
    run.font.bold = bold
    if color_rgb:
        run.font.color.rgb = RGBColor(*color_rgb)
    else:
        run.font.color.rgb = RGBColor(0, 0, 0)

def add_text_box(slide, text, left, top, width, height, font_size=10, bold=False, color_rgb=(0,0,0)):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color_rgb)
    return txBox

def add_bullet_textbox(slide, lines, left, top, width, height, font_size=9, color_rgb=(50,50,50)):
    """Add a text box with bullet lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*color_rgb)
    return txBox

# ─────────────────────────────────────────────
# SLIDE 1: Cover – add title text boxes
# ─────────────────────────────────────────────
s1 = slides[0]
# Add main title
add_text_box(s1,
    "百威商学院 × 柯纳仕\n领导力发展课程方案",
    left=1.5, top=3.5, width=9, height=2.0,
    font_size=28, bold=True, color_rgb=(30, 60, 114))

# ─────────────────────────────────────────────
# SLIDE 2: 百威中国 + 需求理解
# ─────────────────────────────────────────────
s2 = slides[1]
for shape in s2.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        t = shape.text.strip()[:60]
        # Replace the generic intro paragraph
        if '杰出人才是企业成功的基石' in shape.text:
            clear_and_set_text(shape,
                "百威中国是全球最大啤酒酿造商百威英博（AB InBev）亚太区核心市场，"
                "业务覆盖啤酒酿造、进口、推广、分销与销售，运营超过30家酿酒厂，"
                "员工近23,000人，总部位于上海。\n\n"
                "业务的快速扩张对各层级管理者的领导力提出了更高要求。"
                "从销售团队的一线管理者到总部的战略决策层，"
                "都需要系统化的领导力发展体系支撑组织的持续增长。\n\n"
                "柯纳仕基于对快消品行业人才管理挑战的深度理解，"
                "为百威商学院定制覆盖全员管理梯队的领导力发展方案。",
                font_size=12)
        elif '柯纳仕致力于打造' in shape.text:
            clear_and_set_text(shape,
                "百威中国 × 柯纳仕 领导力发展课程方案",
                font_size=20, bold=True)

# ─────────────────────────────────────────────
# SLIDE 3: 我们的使命（mantra – 定制为百威版）
# ─────────────────────────────────────────────
s3 = slides[2]
for shape in s3.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '激发内在动能' in shape.text and len(shape.text) < 100:
            clear_and_set_text(shape,
                "激发内在动能  /  重塑领导行为  /  赋能组织成长",
                font_size=16, bold=True, color_rgb=(30,60,114))
        elif '我们帮助企业识别' in shape.text:
            clear_and_set_text(shape,
                "我们帮助百威中国识别个人、团队、组织与文化的现状，"
                "明确需要到达的位置，通过多样化、体系化的方式激发内在动能，"
                "重塑领导行为，赋能组织成长，支持业务与人才的长效发展。",
                font_size=12)

# ─────────────────────────────────────────────
# SLIDE 4: 我们的服务（保持4大模块标题 + 百威说明）
# ─────────────────────────────────────────────
s4 = slides[3]
for shape in s4.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '我们提供的服务' in shape.text:
            clear_and_set_text(shape, "我们的服务  /  Our Services for Budweiser China", font_size=18, bold=True)

# ─────────────────────────────────────────────
# SLIDE 5: 01 组织诊断 & 变革（保留框架，百威化）
# ─────────────────────────────────────────────
s5 = slides[4]
for shape in s5.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '组织诊断' in shape.text and len(shape.text) < 200:
            clear_and_set_text(shape,
                "01 组织诊断及人才规划\nOrganization Diagnosis & Talent Planning",
                font_size=14, bold=True, color_rgb=(30,60,114))
        elif '组织诊断是企业咨询' in shape.text:
            clear_and_set_text(shape,
                "快速消费品行业竞争激烈，渠道管理、销售团队激励、品牌建设"
                "与供应链效率是企业核心议题。\n\n"
                "柯纳仕为百威中国提供组织诊断服务，综合评估组织能力，"
                "找到人才管理、业务推进与组织协同的切入口，"
                "帮助百威在激烈的市场竞争中将组织能力转化为业绩优势。",
                font_size=11)

# ─────────────────────────────────────────────
# SLIDE 6: 02 人才战略 & 绩效（百威化）
# ─────────────────────────────────────────────
s6 = slides[5]
for shape in s6.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '02 组织人才战略' in shape.text:
            clear_and_set_text(shape,
                "02 人才战略及绩效管理\nTalent Strategy & Performance Management",
                font_size=14, bold=True, color_rgb=(30,60,114))
        elif '业务发展快，人才成长慢' in shape.text:
            clear_and_set_text(shape,
                "百威中国业务发展快，销售与供应链管理人才需求持续增长。"
                "「缺人」成为业务扩张的核心瓶颈——关键岗位人才断层，"
                "外部招募竞争激烈，内部培养体系急需升级。\n\n"
                "我们帮助百威中国打造成功的人才战略，涵盖选用育留，"
                "提升核心销售与运营岗位竞争力，"
                "整合组织绩效与个人绩效，促使商业目标达成！",
                font_size=11)

# ─────────────────────────────────────────────
# SLIDE 7: 03 领导力发展（保留框架）
# ─────────────────────────────────────────────
s7 = slides[6]
for shape in s7.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '03 领导力与高绩效团队发展' in shape.text:
            clear_and_set_text(shape,
                "03 领导力与高绩效团队发展\nLeadership & High-Performance Team Development",
                font_size=14, bold=True, color_rgb=(30,60,114))

# ─────────────────────────────────────────────
# SLIDE 8: 两个维度（保持不变）
# ─────────────────────────────────────────────
# Keep as-is

# ─────────────────────────────────────────────
# SLIDE 9: 纵向自我转型（保持不变）
# ─────────────────────────────────────────────
# Keep as-is

# ─────────────────────────────────────────────
# SLIDE 10: 横向技能扩展 → 替换为详细课程体系
# ─────────────────────────────────────────────
s10 = slides[9]
# Clear and replace with 4-level detailed courses
for shape in s10.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '领导力横向发展' in shape.text:
            clear_and_set_text(shape,
                "课程体系：四层分级 · 精准匹配\nFour-Level Leadership Curriculum for Budweiser China",
                font_size=14, bold=True, color_rgb=(30,60,114))

# Remove existing course list shapes (keep the framework shapes)
shapes_to_remove = []
for shape in s10.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        t = shape.text.strip()
        if ('个人贡献者' in t or '中基层管理者' in t or
            '中高层管理者' in t or 'C-level' in t or
            '完成' in t or '人际' in t or '自我' in t or
            '定目标' in t or '搭团队' in t):
            shapes_to_remove.append(shape)

# Remove course list text boxes
from pptx.oxml import parse_xml
spTree = s10.shapes._spTree
for shape in shapes_to_remove:
    sp = shape._element
    if sp.getparent() is not None:
        sp.getparent().remove(sp)

# Add 4-column course overview
col_data = [
    ("个人贡献者", "完成任务 · 人际协作 · 自我认知",
     ["• 有效沟通（2天）", "• 自信与果敢行为（2天）", "• 非职权影响力（2天）",
      "• 情绪与压力管理（2天）", "• 强有力的演讲（2天）", "• 强有力的工作汇报（2天）",
      "• 问题分析与决策（2天）", "• 结构思考与表达（1/2天）", "• 五代时间管理（2天）",
      "• 项目管理（2天）", "• 销售技巧（2天）", "• 在服务中卓越（2天）",
      "• 冲突管理（2天）", "• DISC/MBTI/优势测评", "• 职业生涯规划"]),
    ("新任经理 / 中基层", "角色转换 · 团队领导 · 绩效达成",
     ["• 管理者的角色认知（2天）", "• 新经理的角色转换与认知（2天）",
      "• 情境领导（2天）", "• 团队领导力（2天）", "• 管理者基本功（2天）",
      "• 绩效管理：提升团队业绩（2天）", "• 授权与激励（2天）",
      "• 教练式管理与辅导（2天）", "• 面试招聘（2天）", "• 建设高绩效团队（2天）",
      "• 高效会议引导（2天）", "• 复盘（1天）", "• 新生代员工管理（2天）",
      "• MBA环：HR for non-HR", "• 非财务人员的财务管理（2天）"]),
    ("中高层管理者", "组织协同 · 业务驱动 · 变革创新",
     ["• 跨部门沟通与协作（2天）", "• 故事影响力（2天）",
      "• 利益相关者管理（2天）", "• 管理者加速器（2天）",
      "• 教练式领导力（2天）", "• 商业敏锐度（2天）", "• 创新思维（2天）",
      "• 变革管理（2天）", "• 建立成长型团队（2天）",
      "• 打造高情商的领导力（2天）", "• 基于神经科学的领导力（2天）",
      "• 打造高绩效团队（2天）"]),
    ("C-level / 高层管理者", "战略解码 · 意识进化 · 组织牵引",
     ["• 领导者意识进化（2天）", "• 内在领导力（2天）",
      "• 人际领导力（2天）", "• 战略解码工作坊（2天）",
      "• 战略人才沙盘（2天）", "• 经营沙盘（2天）",
      "• 变革管理（2天）", "• 利益相关者管理（2天）"]),
]

x_starts = [0.2, 3.4, 6.6, 9.8]
col_w = 2.9
for i, (level, subtitle, items) in enumerate(col_data):
    x = x_starts[i]
    add_text_box(s10, level, x, 1.5, col_w, 0.4, font_size=11, bold=True, color_rgb=(30,60,114))
    add_text_box(s10, subtitle, x, 1.9, col_w, 0.4, font_size=8, color_rgb=(100,100,100))
    add_bullet_textbox(s10, items, x, 2.3, col_w, 4.5, font_size=8, color_rgb=(50,50,50))

# ─────────────────────────────────────────────
# SLIDE 11: 跨文化HR（百威为外资，可保留）
# ─────────────────────────────────────────────
s11 = slides[10]
for shape in s11.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '04 跨文化人力资源' in shape.text:
            clear_and_set_text(shape,
                "04 跨文化人力资源体系管理\nCross Culture HR Management for 百威中国",
                font_size=14, bold=True, color_rgb=(30,60,114))

# ─────────────────────────────────────────────
# SLIDE 12: 详细课程清单附录（替换原有appendix）
# ─────────────────────────────────────────────
s12 = slides[11]
# Clear and create full detailed course listing
for shape in s12.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '附录' in shape.text:
            clear_and_set_text(shape, "附录：柯纳仕详细课程清单  ·  Appendix: Full Course List", font_size=14, bold=True, color_rgb=(30,60,114))
        elif len(shape.text) > 100:
            # Clear the existing giant text block
            clear_and_set_text(shape, "")

# Add course list by level
course_blocks = [
    ("■ 个人贡献者", 5.5,
     ["有效沟通（2天）", "自信与果敢行为（2天）", "非职权影响力（2天）", "情绪与压力管理（2天）",
      "强有力的演讲（2天）", "强有力的工作汇报（2天）", "问题分析与决策（2天）",
      "结构思考与表达（1/2天）", "五代时间管理（2天）", "项目管理（2天）",
      "销售技巧（2天）", "在服务中卓越（2天）", "冲突管理（2天）",
      "DISC性格测试 / MBTI职业性格测试 / 盖洛普五大优势", "职业生涯规划"]),
    ("■ 新任经理 / 中基层管理者", 5.5,
     ["新经理的角色转换与认知（2天）", "管理者的角色认知（2天）", "情境领导（2天）",
      "团队领导力（2天）", "管理者基本功（2天）", "绩效管理：提升团队业绩（2天）",
      "授权与激励（2天）", "教练式管理与辅导（2天）", "面试招聘（2天）",
      "建设高绩效团队（2天）", "高效会议引导（2天）", "复盘（1天）",
      "新生代员工管理（2天）", "非人力资源经理的人力资源管理（2天）",
      "非财务人员的财务管理（2天）", "知己解彼：MBTI与团队领导力（2天）"]),
    ("■ 中高层管理者", 5.5,
     ["跨部门沟通与协作（2天）", "故事影响力（2天）", "利益相关者管理（2天）",
      "管理者加速器（2天）", "教练式领导力（2天）", "商业敏锐度（2天）",
      "创新思维（2天）", "变革管理（2天）", "建立成长型团队（2天）",
      "打造高情商的领导力（2天）", "基于神经科学的领导力（2天）",
      "打造高绩效团队（2天）"]),
    ("■ C-level / 高层管理者", 5.5,
     ["领导者意识进化（2天）", "内在领导力（2天）", "人际领导力（2天）",
      "战略解码工作坊（2天）", "战略人才沙盘（2天）", "经营沙盘（2天）",
      "变革管理（2天）", "利益相关者管理（2天）"]),
]

y_pos = 1.5
for block_title, col_w, items in course_blocks:
    add_text_box(s12, block_title, 0.3, y_pos, col_w, 0.3, font_size=9, bold=True, color_rgb=(30,60,114))
    add_bullet_textbox(s12, items, 0.3, y_pos+0.3, col_w, 5.0, font_size=8, color_rgb=(50,50,50))
    y_pos += 5.2

# ─────────────────────────────────────────────
# SLIDE 13-14: 方法论 + 为什么选择柯纳仕
# ─────────────────────────────────────────────
# Add methodology slide before contact slide
# First clear slide 13 placeholder
s13 = slides[12]
for shape in s13.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '苏州' in shape.text or '杨伟成' in shape.text or '上海' in shape.text or '柯纳仕' in shape.text:
            clear_and_set_text(shape, "")
        elif len(shape.text) > 10:
            clear_and_set_text(shape, "")

# Add methodology title
add_text_box(s13, "方法论体系：让学习真正转化为成果\nMethodology: Turning Learning into Results",
    0.5, 0.3, 11.0, 0.8, font_size=14, bold=True, color_rgb=(30,60,114))

method_cols = [
    ("行动学习\nAction Learning",
     ["基于真实业务问题立项", "结构化反思循环（ORID）",
      "边干边学，解决真实业务挑战", "学习转化与成果交付",
      "已在Brembo等多个项目落地验证"]),
    ("教练辅导\nCoaching",
     ["专业教练一对一辅导", "高管教练与领导力发展",
      "定期跟进与行为追踪", "基于Grow模型的辅导框架",
      "支持从认知到行为的持续转变"]),
    ("复盘机制\nReview",
     ["AAR（After Action Review）", "团队复盘工作坊",
      "个人与组织学习闭环", "将经验转化为组织智慧",
      "推动持续改进的文化落地"]),
    ("测评与觉察\nAssessment",
     ["DISC / MBTI / 盖洛普优势测评", "领导者360度反馈",
      "基于测评数据的自我认知", "为学习路径提供个性化依据",
      "帮助管理者了解自身领导力盲区"]),
]

x_pos = [0.3, 3.1, 6.0, 8.9]
for i, (title, items) in enumerate(method_cols):
    add_text_box(s13, title, x_pos[i], 1.2, 2.5, 0.5, font_size=10, bold=True, color_rgb=(30,60,114))
    add_bullet_textbox(s13, items, x_pos[i], 1.7, 2.5, 3.0, font_size=8, color_rgb=(50,50,50))

# ─────────────────────────────────────────────
# SLIDE 14: 为什么选择柯纳仕（替换contact）
# ─────────────────────────────────────────────
s14 = slides[13]
for shape in s14.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        clear_and_set_text(shape, "")

add_text_box(s14, "为什么选择柯纳仕\nWhy Knas for Budweiser China",
    0.5, 0.2, 11.0, 0.7, font_size=14, bold=True, color_rgb=(30,60,114))

adv_cols = [
    ("三维交付，不是单点课程",
     ["课程解决认知问题", "项目方案解决业务问题", "方法论确保学习转化",
      "测评+教练+复盘+行动学习闭环"]),
    ("四层分级体系，精准匹配",
     ["不是一刀切", "按管理层级精确设计学习路径",
      "覆盖个人贡献者到C-level完整路径"]),
    ("国际标准，本土实践",
     ["顾问团队具备国际背景", "深度理解中国企业实际",
      "服务过BMW、Brembo等标杆客户"]),
    ("结果导向，效果可衡量",
     ["可提供量化ROI与行为改变报告", "长期客户续约率行业领先",
      "17年专注管理培训与组织发展"]),
]

y_pos = 1.1
for i, (title, items) in enumerate(adv_cols):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.2
    y = y_pos + row * 2.3
    add_text_box(s14, f"▶ {title}", x, y, 5.8, 0.4, font_size=10, bold=True, color_rgb=(30,60,114))
    add_bullet_textbox(s14, items, x+0.1, y+0.4, 5.5, 1.8, font_size=9, color_rgb=(50,50,50))

# ─────────────────────────────────────────────
# SLIDE 15: 联系我们（保留）
# ─────────────────────────────────────────────
s15 = slides[14]
# Keep contact info, add Budweiser context
for shape in s15.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        if '团队与领导力' not in shape.text and '联系我们' not in shape.text:
            if len(shape.text) < 50:
                clear_and_set_text(shape, "百威商学院 × 柯纳仕\n携手共创组织成长", font_size=14, bold=True)

# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
prs.save(str(dst))
print(f'Saved: {dst}')
print(f'File size: {dst.stat().st_size / 1024:.1f} KB')
